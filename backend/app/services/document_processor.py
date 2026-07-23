import pdfplumber
from PIL import Image
import pytesseract
import logging
from pptx import Presentation
import docx

logger = logging.getLogger(__name__)

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from a PDF file using pdfplumber."""
    text_content = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text()
            if page_text:
                text_content.append(page_text)
            else:
                logger.warning(f"Could not extract text from page {i} of PDF.")
    return "\n".join(text_content)

def extract_text_from_image(file_path: str) -> str:
    """Extract text from an image using pytesseract OCR."""
    try:
        img = Image.open(file_path)
        return pytesseract.image_to_string(img)
    except Exception as e:
        logger.error(f"Error during image OCR extraction: {str(e)}")
        # Provide a fallback or re-raise
        raise e

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from a modern PowerPoint (.pptx) file."""
    text_content = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        return "\n".join(text_content)
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {str(e)}")
        raise e

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from a Word (.docx) file."""
    text_content = []
    try:
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            text_content.append(para.text)
        return "\n".join(text_content)
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {str(e)}")
        raise e

def extract_text(file_path: str, doc_type: str) -> str:
    """Orchestrate text extraction based on file type."""
    doc_type_lower = doc_type.lower()
    if doc_type_lower == "pdf":
        return extract_text_from_pdf(file_path)
    elif doc_type_lower in ["image", "png", "jpeg", "jpg"]:
        return extract_text_from_image(file_path)
    elif doc_type_lower in ["pptx", "presentation"]:
        return extract_text_from_pptx(file_path)
    elif doc_type_lower in ["docx", "document"]:
        return extract_text_from_docx(file_path)
    elif doc_type_lower == "text":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    else:
        raise ValueError(f"Unsupported document type: {doc_type}")
